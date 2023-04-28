# Creating powerpoint presentations using the python-pptx package
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches

import random
# Create a new PowerPoint presentation
prs = Presentation()

# Add a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Latest Tech in DevOps"
subtitle.text = "Presented by:"

# Define the topics and their corresponding summaries
topics = [
    ("Containerization with Docker", "Containerization makes deploying and scaling applications much easier, with Docker being the most popular containerization platform."),
    ("Continuous Integration with Jenkins", "Jenkins is a popular open-source tool for continuous integration, which helps to automate the build, test and deployment of applications."),
    ("Infrastructure as Code with Terraform", "Terraform is a popular infrastructure as code tool that allows for easy provisioning, configuration and management of infrastructure."),
    ("Monitoring and Logging with Prometheus and ELK Stack", "Prometheus is a popular monitoring and alerting system, while ELK Stack is a popular logging and analysis platform."),
]

# Loop through the topics and create 5 slides for each
for topic, summary in topics:
    # Add a slide with the topic title
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = topic
    
    # Add a summary slide with the topic summary
    summary_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(summary_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Summary"
    body.text = summary
    
    # Add 3 information slides with random content and images
    content_slide_layout = prs.slide_layouts[6]
    for i in range(3):
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = f"{topic} Information Slide {i+1}"
        bullet_points = ["Bullet Point 1", "Bullet Point 2", "Bullet Point 3"]
        for bullet_point in bullet_points:
            p = body.add_paragraph()
            p.text = bullet_point
            p.level = 1
        # Add a random image every 3rd slide
        if i % 3 == 0:
            image_path = f"image_{random.randint(1, 5)}.jpg"
            pic = slide.shapes.add_picture(image_path, Inches(3), Inches(4))
    
# Save the PowerPoint presentation
prs.save("devops.pptx")
